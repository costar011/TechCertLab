import os
import PyPDF2
from pptx import Presentation

def get_all_text(data_folder="data"):
    all_contents = []
    if not os.path.exists(data_folder):
        os.makedirs(data_folder)
        return []

    files = os.listdir(data_folder)
    print(f"📂 폴더 스캔 중: {len(files)}개 파일 발견")

    for file in files:
        path = os.path.join(data_folder, file)
        text = ""
        try:
            if file.lower().endswith('.pdf'):
                with open(path, 'rb') as f:
                    reader = PyPDF2.PdfReader(f)
                    for page in reader.pages:
                        text += page.extract_text() + "\n"
            elif file.lower().endswith('.pptx'):
                prs = Presentation(path)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
            
            if text.strip():
                all_contents.append({"filename": file, "text": text})
                print(f"   ✅ 추출 성공: {file}")
        except Exception as e:
            print(f"   ❌ 에러 발생 ({file}): {e}")

    return all_contents